import streamlit as st
from llm import generate_summary, generate_followup_response
import PyPDF2
from pptx import Presentation
import docx
import chardet


# --- Extraction Functions ---
def extract_text_from_pdf(file_obj):
    """Extract text from a PDF file-like object using PyPDF2."""
    pdf_reader = PyPDF2.PdfReader(file_obj)
    text = ""
    for page in pdf_reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text + "\n"
    return text


def extract_text_from_pptx(file_obj):
    """Extract text from a PPTX file-like object using python-pptx."""
    prs = Presentation(file_obj)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


def extract_text_from_docx(file_obj):
    """Extract text from a DOCX file-like object using python-docx."""
    doc = docx.Document(file_obj)
    text = "\n".join([para.text for para in doc.paragraphs])
    return text


# --- Page 1: Home ---
def home_page():
    st.title("AI Study Assistant Chatbot")
    st.write('Please enter your course name, e.g. "Generative AI"')
    course_name = st.text_input("Course Name")
    if st.button("Submit"):
        if not course_name.strip():
            st.warning("Please enter the course name!")
        else:
            st.session_state['course_name'] = course_name
            st.session_state['page'] = 2
            st.rerun()


# --- Page 2: Upload File ---
def upload_file_page():
    st.title("Upload Files")
    st.write("Please select a file to upload.")

    # Select file category.
    category = st.radio("Select File Category", options=["Reading Materials", "Homework"], index=0)
    st.session_state['file_category'] = category

    file_types = ["pptx", "docx", "pdf", "txt"]
    uploaded_file = st.file_uploader(f"Upload your {category} file", type=file_types, key="file_uploader")

    if uploaded_file is not None:
        ext = uploaded_file.name.split('.')[-1].lower()
        if ext not in file_types:
            st.error("Invalid file type. Please upload a pptx, docx, pdf, or txt file.")
        else:
            if ext == "pdf":
                file_content = extract_text_from_pdf(uploaded_file)
            elif ext == "pptx":
                file_content = extract_text_from_pptx(uploaded_file)
            elif ext == "docx":
                file_content = extract_text_from_docx(uploaded_file)
            elif ext == "txt":
                file_bytes = uploaded_file.read()
                detected = chardet.detect(file_bytes)
                encoding = detected.get("encoding", "utf-8")
                file_content = file_bytes.decode(encoding)
            st.session_state['file_content'] = file_content
            st.session_state['uploaded_file_name'] = uploaded_file.name
            # Set the current file name for tracking summary updates.
            st.session_state['current_file_name'] = uploaded_file.name

            # Maintain global history as a list of file info dictionaries.
            if 'uploaded_files' not in st.session_state:
                st.session_state['uploaded_files'] = []
            st.session_state['uploaded_files'].append({
                "name": uploaded_file.name,
                "category": category,
                "content": file_content,
                "summary": None,       # to be generated later
                "chat_history": []     # per-file conversation history
            })
            # Clear any previous summary from chatbot page.
            st.session_state.pop('initial_output', None)
            st.session_state['page'] = 3
            st.rerun()

    # Show clickable history for previously uploaded files.
    if 'uploaded_files' in st.session_state:
        st.write("**Uploaded Files History:**")
        for idx, file_info in enumerate(st.session_state['uploaded_files']):
            if st.button(f"View Details: {file_info['name']} (Category: {file_info['category']})", key=f"file_{idx}"):
                st.session_state['selected_file_index'] = idx
                st.session_state['page'] = 4
                st.rerun()

    if st.button("Return"):
        st.session_state['page'] = 1
        st.rerun()


# --- Page 3: Chatbot for Current File ---
def chatbot_page():
    st.title("Chatbot - Current File")
    category = st.session_state.get('file_category', 'Unknown')
    file_name = st.session_state.get('uploaded_file_name', 'Unknown File')
    course_name = st.session_state.get('course_name', 'Unknown Course')
    file_content = st.session_state.get('file_content', '')

    st.write(f"**Current File:** {file_name} (Category: {category})")

    # Check if we already have a summary for the current file; if not, process it.
    if ('initial_output' not in st.session_state or
        st.session_state.get('current_file_name') != file_name):
        with st.spinner("Processing current file content..."):
            summary = generate_summary(file_content, category, course_name)
            st.session_state['initial_output'] = summary
            st.session_state['current_file_name'] = file_name
            # Save summary in the corresponding file history.
            for file in st.session_state['uploaded_files']:
                if file['name'] == file_name:
                    file['summary'] = summary

    st.subheader("Initial Output")
    st.write(st.session_state['initial_output'])

    # Build global context from all uploaded files.
    global_context = "File Summaries:\n"
    for file in st.session_state.get('uploaded_files', []):
        if file['name'] == file_name:
            file_summary = st.session_state.get('initial_output', '')
        else:
            if file['summary'] is None:
                file_summary = generate_summary(file['content'], file['category'], course_name)
                file['summary'] = file_summary
            else:
                file_summary = file['summary']
        global_context += f"{file['name']} (Category: {file['category']}):\n{file_summary}\n\n"

    st.write("You can now ask follow-up questions regarding the current file (context is maintained).")
    # Ensure a counter exists for the text input key.
    if 'question_key' not in st.session_state:
        st.session_state['question_key'] = 0

    with st.form(key="followup_form"):
        question_key = f"question_input_{st.session_state['question_key']}"
        user_question = st.text_input("Your Question", key=question_key)
        submitted = st.form_submit_button("Submit Question")

    if submitted and user_question.strip():
        # Build conversation context using your existing code...
        conversation_history = ""
        for file in st.session_state['uploaded_files']:
            if file['name'] == file_name:
                for chat in file['chat_history']:
                    conversation_history += f"Q: {chat['question']}\nA: {chat['answer']}\n"
                break

        full_context = f"Global Context:\n{global_context}\n"
        if conversation_history:
            full_context += f"Conversation History for {file_name}:\n{conversation_history}\n"
        full_context += f"New Question: {user_question}\n"

        response = generate_followup_response(full_context, course_name)
        for file in st.session_state['uploaded_files']:
            if file['name'] == file_name:
                file['chat_history'].append({"question": user_question, "answer": response})
                break

        # Increment the key so that the text input resets.
        st.session_state['question_key'] += 1
        st.rerun()

    st.write("### Chat History for Current File")
    for file in st.session_state.get('uploaded_files', []):
        if file['name'] == file_name:
            for chat in file['chat_history']:
                st.write(f"**Q:** {chat['question']}")
                st.write(f"**A:** {chat['answer']}")
                st.write("---------------------------")
            break

    if st.button("Return"):
        st.session_state['page'] = 2
        st.rerun()


# --- Page 4: Chatbot History ---
def chatbot_history_page():
    st.title("Chatbot - History")
    course_name = st.session_state.get('course_name', 'Unknown Course')
    idx = st.session_state.get('selected_file_index', None)
    if idx is None:
        st.error("No file selected.")
        return
    file = st.session_state['uploaded_files'][idx]
    file_name = file['name']
    st.write(f"**File:** {file_name} (Category: {file['category']})")

    if file['summary'] is None:
        with st.spinner("Processing file content..."):
            file['summary'] = generate_summary(file['content'], file['category'], course_name)
    st.subheader("Initial Output")
    st.write(file['summary'])

    st.write("You can now ask follow-up questions regarding this file (context is maintained).")
    user_question = st.text_input("Your Question", key="detail_question")
    if st.button("Submit Question", key="detail_submit"):
        if user_question.strip():
            conversation_history = ""
            for chat in file.get('chat_history', []):
                conversation_history += f"Q: {chat['question']}\nA: {chat['answer']}\n"
            full_context = f"File Summary:\n{file['summary']}\n"
            if conversation_history:
                full_context += f"Previous Conversation:\n{conversation_history}\n"
            full_context += f"New Question: {user_question}\n"
            response = generate_followup_response(full_context, course_name)
            if 'chat_history' not in file:
                file['chat_history'] = []
            file['chat_history'].append({"question": user_question, "answer": response})
    st.write("### Chat History for This File")
    for chat in file.get('chat_history', []):
        st.write(f"**Q:** {chat['question']}")
        st.write(f"**A:** {chat['answer']}")
        st.write("---------------------------")

    if st.button("Return", key="detail_return"):
        st.session_state['page'] = 2
        st.rerun()


def main():
    if 'page' not in st.session_state:
        st.session_state['page'] = 1

    if st.session_state['page'] == 1:
        home_page()
    elif st.session_state['page'] == 2:
        upload_file_page()
    elif st.session_state['page'] == 3:
        chatbot_page()
    elif st.session_state['page'] == 4:
        chatbot_history_page()


if __name__ == "__main__":
    main()
